import fitz
from docx import Document
from pptx import Presentation
import json
import re
import spacy
from sentence_transformers import SentenceTransformer
import numpy as np
from groq import Groq
from config import GROQ_API_KEY

# Initialize NLP model once at startup
try:
    nlp = spacy.load("en_core_web_sm")
except OSError:
    # Fallback just in case
    import subprocess
    import sys
    subprocess.run([sys.executable, "-m", "spacy", "download", "en_core_web_sm"], check=True)
    nlp = spacy.load("en_core_web_sm")

# Load sentence transformer model locally (runs on CPU)
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

# Initialize Groq client (used exclusively for the dialogue chatbot)
groq_client = Groq(api_key=GROQ_API_KEY)

def extract_text(file_path: str, file_type: str) -> str:
    if file_type == "pdf":
        doc = fitz.open(file_path)
        return "\n".join([page.get_text() for page in doc])
    elif file_type == "docx":
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    elif file_type == "pptx":
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)
    return ""

def generate_talking_points(text: str) -> dict:
    """
    NLP Principle: Extractive Summarization & Dependency Parsing
    Processes the document text with spaCy, scores sentences using word importance (frequencies),
    and transforms the extracted sentences into presentation prompt formats using POS (Part-Of-Speech)
    tagging and linguistic rules.
    Runs 100% locally on CPU.
    """
    if not text.strip():
        return {"title": "Empty Document", "talking_points": []}
    
    # Process text with spaCy (cap at 10000 characters to prevent excessive CPU usage)
    doc = nlp(text[:10000])
    
    # Sentence tokenization
    sentences = list(doc.sents)
    if not sentences:
        return {"title": "Empty Document", "talking_points": []}
        
    # Calculate word frequencies (ignoring stop words, punctuation, and whitespace)
    word_frequencies = {}
    for token in doc:
        if not token.is_stop and not token.is_punct and token.text.strip():
            word = token.text.lower()
            word_frequencies[word] = word_frequencies.get(word, 0) + 1
            
    # Normalize frequencies
    max_freq = max(word_frequencies.values()) if word_frequencies else 1
    for word in word_frequencies:
        word_frequencies[word] = word_frequencies[word] / max_freq
        
    # Score sentences based on the importance of their words
    sentence_scores = {}
    for i, sent in enumerate(sentences):
        score = 0
        word_count = 0
        for token in sent:
            if not token.is_stop and not token.is_punct and token.text.strip():
                word = token.text.lower()
                score += word_frequencies.get(word, 0)
                word_count += 1
        # Normalize score by word count to avoid bias toward longer sentences
        if word_count > 0:
            sentence_scores[i] = score / word_count
        else:
            sentence_scores[i] = 0
            
    # Filter for valid sentences (5 to 45 words long) to keep key points concise and readable
    valid_sentence_indices = [
        i for i, sent in enumerate(sentences) 
        if 5 <= len(sent.text.split()) <= 45
    ]
    
    if not valid_sentence_indices:
        valid_sentence_indices = list(range(len(sentences)))
        
    # Sort indices based on score
    sorted_indices = sorted(valid_sentence_indices, key=lambda x: sentence_scores.get(x, 0), reverse=True)
    
    # Extract top 5 talking points
    top_n = min(len(sorted_indices), 5)
    selected_indices = sorted_indices[:top_n]
    selected_indices.sort() # Sort chronologically
    
    raw_talking_points = [sentences[i] for i in selected_indices]
    formatted_points = []
    
    # Rephrase raw sentences into actionable talking points using spaCy and grammatical rules
    for sent in raw_talking_points:
        tokens = [t for t in sent if not t.is_space]
        start_idx = 0
        # Clean leading transitions
        while start_idx < len(tokens) and (tokens[start_idx].text.lower() in [
            "furthermore", "ultimately", "however", "therefore", "additionally", 
            "moreover", "consequently", "specifically", "indeed", "thus", "hence", 
            "also", "by", "as a result", "for instance", "for example"
        ] or tokens[start_idx].is_punct):
            start_idx += 1
            
        cleaned_tokens = tokens[start_idx:]
        if not cleaned_tokens:
            continue
            
        # Reconstruct sentence text
        cleaned_text = "".join([t.text_with_ws for t in cleaned_tokens]).strip()
        cleaned_text = cleaned_text[0].upper() + cleaned_text[1:]
        if cleaned_text.endswith("."):
            cleaned_text = cleaned_text[:-1]
            
        first_token = cleaned_tokens[0]
        first_char_lower = cleaned_text[0].lower() + cleaned_text[1:]
        cleaned_lower = cleaned_text.lower()
        
        # Rule-based rephrasing utilizing spaCy POS Tags to format as presentation cues
        if first_token.pos_ == "VERB" or first_token.tag_ == "VBG":  # Starts with a verb/gerund, e.g., "Practicing..."
            prompt = f"Discuss the role of {first_char_lower}."
        elif "help" in cleaned_lower or "helps" in cleaned_lower or "benefit" in cleaned_lower or "enhance" in cleaned_lower or "improve" in cleaned_lower:
            prompt = f"Explain how {first_char_lower}."
        elif "important" in cleaned_lower or "critical" in cleaned_lower or "essential" in cleaned_lower or "key" in cleaned_lower:
            prompt = f"Discuss the importance of how {first_char_lower}."
        else:
            prompt = f"Explain why {first_char_lower}."
            
        formatted_points.append(prompt)
    
    # Generate a brief title: first line (if short) or the most frequent noun chunk
    title = ""
    first_line = text.split("\n")[0].strip()
    if first_line and len(first_line) < 50:
        title = first_line
    else:
        # Fallback: most common noun chunk
        noun_chunks = [chunk.text.strip().title() for chunk in doc.noun_chunks if len(chunk.text.split()) < 4]
        if noun_chunks:
            chunk_counts = {}
            for chunk in noun_chunks:
                chunk_counts[chunk] = chunk_counts.get(chunk, 0) + 1
            title = max(chunk_counts, key=chunk_counts.get)
        else:
            title = "Extracted Document Summary"
            
    return {
        "title": title,
        "talking_points": formatted_points
    }

def check_relevance(topic: str, transcript: str) -> dict:
    """
    NLP Principle: Sentence Embeddings & Cosine Similarity
    Encodes the topic and user transcript into vector representations, computes their similarity,
    and returns a score mapped to 0-100 with rule-based feedback on vocabulary overlap.
    """
    if not topic.strip() or not transcript.strip():
        return {
            "relevance_score": 0,
            "relevance_feedback": "No topic or transcript available to evaluate."
        }
        
    # Generate vector embeddings locally
    vectors = embedding_model.encode([topic, transcript])
    vec_topic, vec_transcript = vectors[0], vectors[1]
    
    # Compute Cosine Similarity
    dot_product = np.dot(vec_topic, vec_transcript)
    norm_topic = np.linalg.norm(vec_topic)
    norm_transcript = np.linalg.norm(vec_transcript)
    similarity = dot_product / (norm_topic * norm_transcript) if (norm_topic * norm_transcript) > 0 else 0.0
    
    # Map similarity (typically in 0.1 to 0.7 range) to a 0-100 score
    score = int(max(0.0, min(1.0, (similarity - 0.1) / 0.55)) * 100)
    
    # Extract key vocabulary overlapping between topic and transcript for feedback
    doc_topic = nlp(topic.lower())
    doc_transcript = nlp(transcript.lower())
    
    def get_keywords(doc):
        return {token.lemma_ for token in doc if token.pos_ in ("NOUN", "PROPN", "VERB") and not token.is_stop}
        
    words_topic = get_keywords(doc_topic)
    words_transcript = get_keywords(doc_transcript)
    
    overlap = words_topic.intersection(words_transcript)
    missed = words_topic.difference(words_transcript)
    
    if score >= 75:
        feedback = f"Excellent relevance! You successfully addressed the topic and directly spoke about key concepts like: {', '.join(list(overlap)[:3])}."
    elif score >= 50:
        feedback = f"Good focus, but you could align closer. You mentioned {', '.join(list(overlap)[:2])}, but try to elaborate more on {', '.join(list(missed)[:2])}."
    else:
        feedback = f"Low relevance. The speech drifted from the main topic. Ensure you define and speak directly about: {', '.join(list(words_topic)[:3])}."
        
    return {
        "relevance_score": score,
        "relevance_feedback": feedback
    }

def check_coverage(talking_points: list, transcript: str) -> dict:
    """
    NLP Principle: Semantic Matching via Vector Embeddings
    Segments the speech transcript into sentences, embeds them alongside the target talking points,
    and runs a similarity search to check which talking points were covered.
    """
    if not talking_points or not transcript.strip():
        return {
            "coverage_report": [],
            "coverage_score": 0,
            "coverage_feedback": "No talking points or transcript available to evaluate."
        }
        
    # Split transcript into sentences
    doc = nlp(transcript)
    transcript_sentences = [sent.text.strip() for sent in doc.sents if len(sent.text.split()) > 2]
    
    if not transcript_sentences:
        transcript_sentences = [transcript]
        
    # Embed talking points and transcript sentences
    tp_embeddings = embedding_model.encode(talking_points)
    sent_embeddings = embedding_model.encode(transcript_sentences)
    
    # Calculate cosine similarity matrix
    tp_norms = np.linalg.norm(tp_embeddings, axis=1, keepdims=True)
    sent_norms = np.linalg.norm(sent_embeddings, axis=1, keepdims=True)
    similarity_matrix = np.dot(tp_embeddings, sent_embeddings.T) / (np.dot(tp_norms, sent_norms.T) + 1e-9)
    
    # Similarity threshold to match a talking point
    threshold = 0.55
    
    coverage_report = []
    covered_count = 0
    
    for i, tp in enumerate(talking_points):
        max_idx = np.argmax(similarity_matrix[i])
        max_sim = similarity_matrix[i][max_idx]
        
        covered = bool(max_sim >= threshold)
        confidence = "low"
        if max_sim >= 0.70:
            confidence = "high"
        elif max_sim >= 0.55:
            confidence = "medium"
            
        if covered:
            covered_count += 1
            feedback = f"Matched with high similarity to: \"{transcript_sentences[max_idx]}\""
        else:
            feedback = "This point was not covered. Try to explain this concept clearly in your speech."
            
        coverage_report.append({
            "talking_point": tp,
            "covered": covered,
            "confidence": confidence,
            "feedback": feedback
        })
        
    coverage_score = int((covered_count / len(talking_points)) * 100)
    
    if coverage_score >= 80:
        feedback_summary = "Excellent content coverage! You verbally explained almost all of your talking points."
    elif coverage_score >= 50:
        feedback_summary = "Moderate coverage. You hit some key points, but skipped several important areas."
    else:
        feedback_summary = "Low coverage. You missed most of the prepared talking points in your speech."
        
    return {
        "coverage_report": coverage_report,
        "coverage_score": coverage_score,
        "coverage_feedback": feedback_summary
    }

def chat_with_coach(history: list, system_prompt: str) -> str:
    """
    Dialogue Agent: Groq Cloud API
    Uses Llama-3-8B to generate coach advice for public speaking queries.
    """
    if not GROQ_API_KEY:
        return "Coach: API Key Missing. Please set your GROQ_API_KEY in config.py to talk to the speech coach."
        
    # Convert conversational history to OpenAI/Groq format
    messages = [{"role": "system", "content": system_prompt}]
    for msg in history:
        role = "user" if msg["role"] == "user" else "assistant"
        messages.append({"role": role, "content": msg["content"]})
        
    try:
        chat_completion = groq_client.chat.completions.create(
            messages=messages,
            model="llama-3.1-8b-instant",
            max_tokens=500,
        )
        return chat_completion.choices[0].message.content.strip()
    except Exception as e:
        return f"Coach: Error communicating with Groq API ({str(e)}). Please verify your GROQ_API_KEY in config.py."